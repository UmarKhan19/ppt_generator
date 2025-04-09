from flask import Flask, request, send_file
from pptx import Presentation
from pptx.util import Pt
from pptx.enum.shapes import MSO_SHAPE_TYPE
import tempfile
import os
import json
import copy
from io import BytesIO

app = Flask(__name__)

class PresentationBuilder:
    def __init__(self, template_path):
        print(f"[INIT] Loading template from: {template_path}")
        self.prs = Presentation(template_path)
        self.base_slide = self.prs.slides[0]  # assuming first slide is clean university template
        print("[INIT] Base slide loaded")

    def duplicate_base_slide(self):
        print("[SLIDE] Duplicating base slide")
        slide_layout = self.base_slide.slide_layout
        new_slide = self.prs.slides.add_slide(slide_layout)
        for shape in self.base_slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                print("[WARN] Skipping grouped shape")
                continue
            if shape.is_placeholder:
                continue
            el = shape.element
            new_slide.shapes._spTree.insert_element_before(copy.deepcopy(el), 'p:extLst')
        print("[SLIDE] Slide duplicated")
        return new_slide

    def add_content_slide(self, title, content):
        print(f"[CONTENT] Adding slide: {title}")
        slide = self.duplicate_base_slide()
        title_placeholder = None
        content_placeholder = None
        for shape in slide.shapes:
            if not shape.is_placeholder:
                continue
            if shape.placeholder_format.idx == 0:
                title_placeholder = shape
            else:
                content_placeholder = shape
        if title_placeholder:
            truncated_title = title if len(title) < 200 else title[:197] + '...'
            title_placeholder.text = truncated_title
        if content_placeholder:
            content_frame = content_placeholder.text_frame
            content_frame.clear()
            for line in content.strip().split('\n'):
                clean_line = line.strip()
                if not clean_line:
                    continue
                p = content_frame.add_paragraph()
                p.text = clean_line
                p.font.size = Pt(18)
                p.level = 0
        print(f"[CONTENT] Slide '{title}' added successfully")

    def build_from_content(self, content_dict):
        print("[BUILD] Building presentation from content")
        for section, slides in content_dict.items():
            print(f"[SECTION] Processing: {section}")
            self.add_content_slide(section, "")
            for slide_data in slides:
                try:
                    self.add_content_slide(slide_data['title'], slide_data['content'])
                except KeyError as e:
                    print(f"[ERROR] Missing key in slide data: {e}")
        print("[BUILD] Presentation build complete")

    def save(self, output_path):
        print(f"[SAVE] Saving presentation to: {output_path}")
        self.prs.save(output_path)
        print("[SAVE] Save complete")

@app.route('/generate-ppt', methods=['POST'])
def generate_ppt():
    print("[REQUEST] /generate-ppt POST received")
    template = request.files.get('template')
    content_file = request.files.get('content')
    if not template or not content_file:
        print("[ERROR] Missing files in request")
        return {"error": "Both template and content files are required."}, 400

    with tempfile.TemporaryDirectory() as tmpdir:
        print(f"[TEMP] Temporary directory created at: {tmpdir}")
        template_path = os.path.join(tmpdir, 'template.pptx')
        content_path = os.path.join(tmpdir, 'content.json')
        output_path = os.path.join(tmpdir, 'output.pptx')

        print("[FILE] Saving uploaded files")
        template.save(template_path)
        content_file.save(content_path)

        print("[FILE] Reading content JSON")
        try:
            with open(content_path, 'r', encoding='utf-8') as f:
                content = json.load(f)
        except json.JSONDecodeError as e:
            print(f"[ERROR] JSON decode failed: {e}")
            return {"error": "Invalid JSON content file."}, 400

        print("[PROCESS] Generating presentation")
        builder = PresentationBuilder(template_path)
        builder.build_from_content(content)
        builder.save(output_path)

        print("[RESPONSE] Sending file to client")
        return send_file(output_path, as_attachment=True, download_name='Generated_Presentation.pptx')

@app.route('/health', methods=['GET'])
def health_check():
    print("[HEALTH] Health check request received")
    return {"status": "Service is healthy"}, 200

if __name__ == '__main__':
    print("[SERVER] Starting Flask app on port 5000")
    app.run(debug=True, port=5000)
